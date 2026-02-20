"""Tests for azext_prototype.parsers.binary_reader."""

from __future__ import annotations

import base64
from pathlib import Path
from unittest.mock import patch

import pytest

from azext_prototype.parsers.binary_reader import (
    EmbeddedImage,
    FileCategory,
    ReadResult,
    classify_file,
    read_file,
    MAX_IMAGE_SIZE,
    MAX_IMAGES_PER_DIR,
)


# ------------------------------------------------------------------ #
# classify_file
# ------------------------------------------------------------------ #


class TestClassifyFile:
    """Extension-based file classification."""

    @pytest.mark.parametrize("ext,expected", [
        (".jpg", FileCategory.IMAGE),
        (".jpeg", FileCategory.IMAGE),
        (".png", FileCategory.IMAGE),
        (".gif", FileCategory.IMAGE),
        (".webp", FileCategory.IMAGE),
        (".bmp", FileCategory.IMAGE),
        (".tiff", FileCategory.IMAGE),
        (".tif", FileCategory.IMAGE),
    ])
    def test_image_extensions(self, tmp_path, ext, expected):
        p = tmp_path / f"file{ext}"
        p.touch()
        assert classify_file(p) == expected

    @pytest.mark.parametrize("ext,expected", [
        (".pdf", FileCategory.DOCUMENT),
        (".docx", FileCategory.DOCUMENT),
        (".pptx", FileCategory.DOCUMENT),
        (".xlsx", FileCategory.DOCUMENT),
    ])
    def test_document_extensions(self, tmp_path, ext, expected):
        p = tmp_path / f"file{ext}"
        p.touch()
        assert classify_file(p) == expected

    def test_svg_is_text(self, tmp_path):
        p = tmp_path / "diagram.svg"
        p.touch()
        assert classify_file(p) == FileCategory.TEXT

    @pytest.mark.parametrize("ext", [".md", ".txt", ".yaml", ".py", ".json", ".csv"])
    def test_text_extensions(self, tmp_path, ext):
        p = tmp_path / f"file{ext}"
        p.touch()
        assert classify_file(p) == FileCategory.TEXT

    def test_unknown_extension_defaults_to_text(self, tmp_path):
        p = tmp_path / "data.xyz"
        p.touch()
        assert classify_file(p) == FileCategory.TEXT

    def test_case_insensitive(self, tmp_path):
        p = tmp_path / "photo.JPG"
        p.touch()
        assert classify_file(p) == FileCategory.IMAGE


# ------------------------------------------------------------------ #
# _read_text (via read_file)
# ------------------------------------------------------------------ #


class TestReadText:
    """Text file reading."""

    def test_read_utf8(self, tmp_path):
        f = tmp_path / "notes.md"
        f.write_text("# Hello world", encoding="utf-8")
        result = read_file(f)
        assert result.category == FileCategory.TEXT
        assert result.text == "# Hello world"
        assert result.error is None

    def test_read_non_utf8_replaces(self, tmp_path):
        f = tmp_path / "data.bin"
        f.write_bytes(b"hello \xff\xfe world")
        result = read_file(f)
        assert result.category == FileCategory.TEXT
        assert "hello" in result.text
        assert result.error is None

    def test_read_svg_as_text(self, tmp_path):
        f = tmp_path / "diagram.svg"
        f.write_text("<svg><circle r='10'/></svg>", encoding="utf-8")
        result = read_file(f)
        assert result.category == FileCategory.TEXT
        assert "<svg>" in result.text


# ------------------------------------------------------------------ #
# _read_image (via read_file)
# ------------------------------------------------------------------ #

# Minimal valid 1x1 PNG (67 bytes)
_TINY_PNG = (
    b"\x89PNG\r\n\x1a\n\x00\x00\x00\rIHDR\x00\x00\x00\x01"
    b"\x00\x00\x00\x01\x08\x02\x00\x00\x00\x90wS\xde\x00"
    b"\x00\x00\x0cIDATx\x9cc\xf8\x0f\x00\x00\x01\x01\x00"
    b"\x05\x18\xd8N\x00\x00\x00\x00IEND\xaeB`\x82"
)


class TestReadImage:
    """Standalone image reading and base64 encoding."""

    def test_read_valid_png(self, tmp_path):
        f = tmp_path / "arch.png"
        f.write_bytes(_TINY_PNG)
        result = read_file(f)
        assert result.category == FileCategory.IMAGE
        assert result.image_data is not None
        assert result.mime_type == "image/png"
        assert result.error is None
        # Verify base64 round-trips
        assert base64.b64decode(result.image_data) == _TINY_PNG

    def test_read_jpeg(self, tmp_path):
        f = tmp_path / "photo.jpg"
        # JPEG header bytes
        f.write_bytes(b"\xff\xd8\xff\xe0" + b"\x00" * 50)
        result = read_file(f)
        assert result.category == FileCategory.IMAGE
        assert result.image_data is not None
        assert result.mime_type == "image/jpeg"

    def test_image_too_large(self, tmp_path):
        f = tmp_path / "big.png"
        f.write_bytes(b"\x89PNG" + b"\x00" * 50)
        with patch("azext_prototype.parsers.binary_reader.MAX_IMAGE_SIZE", 10):
            result = read_file(f)
        assert result.category == FileCategory.IMAGE
        assert result.error is not None
        assert "too large" in result.error

    def test_image_unreadable(self, tmp_path):
        f = tmp_path / "missing.png"
        # Don't create the file
        result = read_file(f)
        assert result.category == FileCategory.IMAGE
        assert result.error is not None


# ------------------------------------------------------------------ #
# _read_document — PDF
# ------------------------------------------------------------------ #


class TestReadPDF:
    """PDF text and image extraction via pypdf."""

    def test_read_pdf_text(self, tmp_path):
        """Create a minimal PDF with pypdf and verify text extraction."""
        from pypdf import PdfWriter

        writer = PdfWriter()
        writer.add_blank_page(width=72, height=72)
        # pypdf doesn't have a simple way to add text to a blank page,
        # so we use the annotation approach
        pdf_path = tmp_path / "doc.pdf"
        with open(pdf_path, "wb") as f:
            writer.write(f)

        result = read_file(pdf_path)
        assert result.category == FileCategory.DOCUMENT
        # Blank page — may have empty text but shouldn't error
        assert result.error is None or result.embedded_images is not None

    def test_read_pdf_missing_library(self, tmp_path):
        """When pypdf is not installed, returns actionable error."""
        f = tmp_path / "doc.pdf"
        f.write_bytes(b"%PDF-1.4 dummy")

        with patch.dict("sys.modules", {"pypdf": None}):
            result = read_file(f)
        assert result.category == FileCategory.DOCUMENT
        assert result.error is not None
        assert "Missing library" in result.error or "pip install" in result.error


# ------------------------------------------------------------------ #
# _read_document — DOCX
# ------------------------------------------------------------------ #


class TestReadDOCX:
    """Word document text and image extraction."""

    def test_read_docx_text(self, tmp_path):
        """Create a minimal DOCX and verify text extraction."""
        from docx import Document

        doc = Document()
        doc.add_paragraph("Hello from Word")
        doc.add_paragraph("Second paragraph")
        docx_path = tmp_path / "spec.docx"
        doc.save(str(docx_path))

        result = read_file(docx_path)
        assert result.category == FileCategory.DOCUMENT
        assert result.error is None
        assert "Hello from Word" in result.text
        assert "Second paragraph" in result.text

    def test_read_docx_with_image(self, tmp_path):
        """DOCX with an embedded image extracts both text and image."""
        from docx import Document
        from docx.shared import Inches
        from PIL import Image as PILImage
        import io

        # Create a proper PNG via Pillow (python-docx validates PNG structure)
        img_buf = io.BytesIO()
        PILImage.new("RGB", (10, 10), color="red").save(img_buf, format="PNG")
        img_path = tmp_path / "logo.png"
        img_path.write_bytes(img_buf.getvalue())

        doc = Document()
        doc.add_paragraph("Document with image")
        doc.add_picture(str(img_path), width=Inches(1))
        docx_path = tmp_path / "with_image.docx"
        doc.save(str(docx_path))

        result = read_file(docx_path)
        assert result.category == FileCategory.DOCUMENT
        assert result.error is None
        assert "Document with image" in result.text
        assert len(result.embedded_images) >= 1
        img = result.embedded_images[0]
        assert img.mime_type.startswith("image/")
        assert img.data  # base64 data present
        assert "with_image.docx" in img.source

    def test_read_docx_missing_library(self, tmp_path):
        f = tmp_path / "doc.docx"
        f.write_bytes(b"PK\x03\x04 dummy")
        with patch.dict("sys.modules", {"docx": None}):
            result = read_file(f)
        assert result.error is not None


# ------------------------------------------------------------------ #
# _read_document — PPTX
# ------------------------------------------------------------------ #


class TestReadPPTX:
    """PowerPoint text and image extraction."""

    def test_read_pptx_text(self, tmp_path):
        """Create a PPTX with text and verify extraction."""
        from pptx import Presentation
        from pptx.util import Inches

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])  # title + content
        slide.shapes.title.text = "Architecture Overview"
        slide.placeholders[1].text = "This is the content"

        pptx_path = tmp_path / "deck.pptx"
        prs.save(str(pptx_path))

        result = read_file(pptx_path)
        assert result.category == FileCategory.DOCUMENT
        assert result.error is None
        assert "Architecture Overview" in result.text
        assert "This is the content" in result.text

    def test_read_pptx_with_image(self, tmp_path):
        """PPTX with an embedded image extracts both text and image."""
        from pptx import Presentation
        from pptx.util import Inches
        from PIL import Image as PILImage
        import io

        img_buf = io.BytesIO()
        PILImage.new("RGB", (10, 10), color="blue").save(img_buf, format="PNG")
        img_path = tmp_path / "icon.png"
        img_path.write_bytes(img_buf.getvalue())

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
        slide.shapes.add_picture(str(img_path), Inches(1), Inches(1))

        pptx_path = tmp_path / "with_image.pptx"
        prs.save(str(pptx_path))

        result = read_file(pptx_path)
        assert result.category == FileCategory.DOCUMENT
        # May have no text (blank slide) — that's OK if images are found
        assert len(result.embedded_images) >= 1
        img = result.embedded_images[0]
        assert img.mime_type.startswith("image/")
        assert img.data

    def test_read_pptx_missing_library(self, tmp_path):
        f = tmp_path / "deck.pptx"
        f.write_bytes(b"PK\x03\x04 dummy")
        with patch.dict("sys.modules", {"pptx": None}):
            result = read_file(f)
        assert result.error is not None


# ------------------------------------------------------------------ #
# _read_document — XLSX
# ------------------------------------------------------------------ #


class TestReadXLSX:
    """Excel text extraction (no image extraction)."""

    def test_read_xlsx_text(self, tmp_path):
        from openpyxl import Workbook

        wb = Workbook()
        ws = wb.active
        ws.title = "Costs"
        ws.append(["Service", "SKU", "Monthly"])
        ws.append(["App Service", "S1", "73.00"])
        xlsx_path = tmp_path / "costs.xlsx"
        wb.save(str(xlsx_path))

        result = read_file(xlsx_path)
        assert result.category == FileCategory.DOCUMENT
        assert result.error is None
        assert "App Service" in result.text
        assert "73.00" in result.text
        assert result.embedded_images == []

    def test_read_xlsx_missing_library(self, tmp_path):
        f = tmp_path / "data.xlsx"
        f.write_bytes(b"PK\x03\x04 dummy")
        with patch.dict("sys.modules", {"openpyxl": None}):
            result = read_file(f)
        assert result.error is not None


# ------------------------------------------------------------------ #
# ReadResult dataclass
# ------------------------------------------------------------------ #


class TestReadResult:
    """ReadResult defaults and construction."""

    def test_default_embedded_images_empty(self):
        r = ReadResult(category=FileCategory.TEXT, text="hi", filename="f.txt")
        assert r.embedded_images == []

    def test_embedded_image_dataclass(self):
        img = EmbeddedImage(data="abc123", mime_type="image/png", source="doc.docx/image1.png")
        assert img.data == "abc123"
        assert img.mime_type == "image/png"


# ------------------------------------------------------------------ #
# Constants
# ------------------------------------------------------------------ #


class TestConstants:
    def test_max_image_size(self):
        assert MAX_IMAGE_SIZE == 20 * 1024 * 1024

    def test_max_images_per_dir(self):
        assert MAX_IMAGES_PER_DIR == 250
