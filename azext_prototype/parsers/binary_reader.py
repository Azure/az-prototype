"""Binary file reading — type detection, text extraction, and image encoding.

Classifies files into three categories:
  - TEXT: Read as UTF-8 (existing behavior)
  - DOCUMENT: Extract text content + embedded images (PDF, DOCX, PPTX, XLSX)
  - IMAGE: Base64-encode for vision API (JPG, PNG, GIF, WebP, BMP, TIFF)

SVG files are classified as TEXT (they are XML).
"""

from __future__ import annotations

import base64
import logging
import mimetypes
from dataclasses import dataclass, field
from enum import Enum
from pathlib import Path

logger = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# Constants
# ---------------------------------------------------------------------------

MAX_IMAGE_SIZE = 20 * 1024 * 1024  # 20 MB per image
MAX_IMAGES_PER_DIR = 250  # accommodates large documents with many embedded images

# ---------------------------------------------------------------------------
# Data structures
# ---------------------------------------------------------------------------


class FileCategory(Enum):
    TEXT = "text"
    DOCUMENT = "document"
    IMAGE = "image"


@dataclass
class EmbeddedImage:
    """An image extracted from inside a document (PDF, DOCX, PPTX)."""

    data: str  # base64-encoded image bytes
    mime_type: str  # e.g. "image/png"
    source: str  # human-readable origin, e.g. "requirements.docx/image1.jpeg"


@dataclass
class ReadResult:
    """Result of reading a single file."""

    category: FileCategory
    text: str | None = None
    image_data: str | None = None  # base64-encoded (standalone images only)
    mime_type: str | None = None  # MIME type (standalone images only)
    filename: str = ""
    error: str | None = None
    embedded_images: list[EmbeddedImage] = field(default_factory=list)


# ---------------------------------------------------------------------------
# Extension → category mapping
# ---------------------------------------------------------------------------

_IMAGE_EXTENSIONS = frozenset({
    ".jpg", ".jpeg", ".png", ".gif", ".webp", ".bmp", ".tiff", ".tif",
})
_DOCUMENT_EXTENSIONS = frozenset({".pdf", ".docx", ".pptx", ".xlsx"})
_TEXT_EXTENSIONS_OVERRIDE = frozenset({".svg"})  # SVG is XML text


# ---------------------------------------------------------------------------
# Public API
# ---------------------------------------------------------------------------


def classify_file(path: Path) -> FileCategory:
    """Classify a file by its extension."""
    ext = path.suffix.lower()
    if ext in _TEXT_EXTENSIONS_OVERRIDE:
        return FileCategory.TEXT
    if ext in _IMAGE_EXTENSIONS:
        return FileCategory.IMAGE
    if ext in _DOCUMENT_EXTENSIONS:
        return FileCategory.DOCUMENT
    return FileCategory.TEXT  # default: attempt text read


def read_file(path: Path) -> ReadResult:
    """Read a file, dispatching to the appropriate handler by category."""
    category = classify_file(path)
    filename = path.name

    if category == FileCategory.IMAGE:
        return _read_image(path, filename)
    elif category == FileCategory.DOCUMENT:
        return _read_document(path, filename)
    else:
        return _read_text(path, filename)


# ---------------------------------------------------------------------------
# Text files
# ---------------------------------------------------------------------------


def _read_text(path: Path, filename: str) -> ReadResult:
    """Read file as UTF-8 text."""
    try:
        text = path.read_text(encoding="utf-8", errors="replace")
        return ReadResult(category=FileCategory.TEXT, text=text, filename=filename)
    except Exception as e:
        return ReadResult(category=FileCategory.TEXT, filename=filename, error=str(e))


# ---------------------------------------------------------------------------
# Standalone images
# ---------------------------------------------------------------------------


def _read_image(path: Path, filename: str) -> ReadResult:
    """Base64-encode an image file for the vision API."""
    try:
        size = path.stat().st_size
        if size > MAX_IMAGE_SIZE:
            return ReadResult(
                category=FileCategory.IMAGE,
                filename=filename,
                error=f"Image too large ({size // (1024 * 1024)}MB > {MAX_IMAGE_SIZE // (1024 * 1024)}MB limit)",
            )
        data = base64.b64encode(path.read_bytes()).decode("utf-8")
        mime = mimetypes.guess_type(str(path))[0] or "image/png"
        return ReadResult(
            category=FileCategory.IMAGE,
            image_data=data,
            mime_type=mime,
            filename=filename,
        )
    except Exception as e:
        return ReadResult(category=FileCategory.IMAGE, filename=filename, error=str(e))


# ---------------------------------------------------------------------------
# Documents (text + embedded images)
# ---------------------------------------------------------------------------


def _read_document(path: Path, filename: str) -> ReadResult:
    """Extract text and embedded images from a document file."""
    ext = path.suffix.lower()
    try:
        if ext == ".pdf":
            text, images = _extract_pdf(path, filename)
        elif ext == ".docx":
            text, images = _extract_docx(path, filename)
        elif ext == ".pptx":
            text, images = _extract_pptx(path, filename)
        elif ext == ".xlsx":
            text, images = _extract_xlsx(path, filename)
        else:
            return ReadResult(
                category=FileCategory.DOCUMENT,
                filename=filename,
                error=f"Unsupported document type: {ext}",
            )

        if not text or not text.strip():
            # Still return embedded images even if no text was extracted
            if images:
                return ReadResult(
                    category=FileCategory.DOCUMENT,
                    text="[No text content extracted]",
                    filename=filename,
                    embedded_images=images,
                )
            return ReadResult(
                category=FileCategory.DOCUMENT,
                filename=filename,
                error="No text content could be extracted",
            )

        return ReadResult(
            category=FileCategory.DOCUMENT,
            text=text,
            filename=filename,
            embedded_images=images,
        )
    except ImportError as e:
        logger.warning("Missing library for %s: %s", ext, e)
        return ReadResult(
            category=FileCategory.DOCUMENT,
            filename=filename,
            error=f"Missing library: {e}. Install with: pip install pypdf python-docx python-pptx openpyxl",
        )
    except Exception as e:
        logger.warning("Failed to extract %s: %s", filename, e)
        return ReadResult(
            category=FileCategory.DOCUMENT,
            filename=filename,
            error=str(e),
        )


# ---------------------------------------------------------------------------
# Format-specific extractors
# ---------------------------------------------------------------------------


def _encode_blob(blob: bytes) -> str:
    """Base64-encode raw image bytes."""
    return base64.b64encode(blob).decode("utf-8")


def _extract_pdf(path: Path, filename: str) -> tuple[str, list[EmbeddedImage]]:
    """Extract text and embedded images from a PDF."""
    from pypdf import PdfReader

    reader = PdfReader(str(path))
    pages: list[str] = []
    images: list[EmbeddedImage] = []

    for i, page in enumerate(reader.pages):
        text = page.extract_text()
        if text:
            pages.append(f"[Page {i + 1}]\n{text}")

        # Extract embedded images
        try:
            for img_key in page.images.keys():
                img = page.images[img_key]
                blob = img.data
                if len(blob) > MAX_IMAGE_SIZE:
                    continue
                mime = mimetypes.guess_type(img.name)[0] or "image/png"
                images.append(EmbeddedImage(
                    data=_encode_blob(blob),
                    mime_type=mime,
                    source=f"{filename}/page{i + 1}/{img.name}",
                ))
        except Exception as e:
            logger.debug("Could not extract images from PDF page %d: %s", i + 1, e)

    return "\n\n".join(pages), images


def _extract_docx(path: Path, filename: str) -> tuple[str, list[EmbeddedImage]]:
    """Extract text and embedded images from a Word document."""
    from docx import Document

    doc = Document(str(path))

    # Extract text
    text = "\n\n".join(p.text for p in doc.paragraphs if p.text.strip())

    # Extract embedded images via relationship parts
    images: list[EmbeddedImage] = []
    try:
        for rel in doc.part.rels.values():
            if "image" in rel.reltype:
                try:
                    blob = rel.target_part.blob
                    if len(blob) > MAX_IMAGE_SIZE:
                        continue
                    mime = rel.target_part.content_type
                    name = rel.target_ref if hasattr(rel, "target_ref") else "image"
                    images.append(EmbeddedImage(
                        data=_encode_blob(blob),
                        mime_type=mime,
                        source=f"{filename}/{name}",
                    ))
                except Exception as e:
                    logger.debug("Could not extract image from DOCX rel: %s", e)
    except Exception as e:
        logger.debug("Could not enumerate DOCX image rels: %s", e)

    return text, images


def _extract_pptx(path: Path, filename: str) -> tuple[str, list[EmbeddedImage]]:
    """Extract text and embedded images from a PowerPoint presentation."""
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(str(path))
    slides: list[str] = []
    images: list[EmbeddedImage] = []

    for i, slide in enumerate(prs.slides):
        texts: list[str] = []
        for shape in slide.shapes:
            # Extract text from text frames
            if shape.has_text_frame:
                frame_text = shape.text_frame.text
                if frame_text.strip():
                    texts.append(frame_text)

            # Extract images from picture shapes
            try:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    blob = shape.image.blob
                    if len(blob) > MAX_IMAGE_SIZE:
                        continue
                    mime = shape.image.content_type
                    ext = shape.image.ext
                    images.append(EmbeddedImage(
                        data=_encode_blob(blob),
                        mime_type=mime,
                        source=f"{filename}/slide{i + 1}.{ext}",
                    ))
            except (AttributeError, Exception) as e:
                # Known issue: JPEG extraction can fail on some files
                logger.debug("Could not extract image from PPTX slide %d: %s", i + 1, e)

        if texts:
            slides.append(f"[Slide {i + 1}]\n" + "\n".join(texts))

    return "\n\n".join(slides), images


def _extract_xlsx(path: Path, filename: str) -> tuple[str, list[EmbeddedImage]]:
    """Extract text from an Excel workbook (no image extraction)."""
    from openpyxl import load_workbook

    wb = load_workbook(str(path), read_only=True, data_only=True)
    sheets: list[str] = []

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows: list[str] = []
        for row in ws.iter_rows(values_only=True):
            row_text = " | ".join(str(c) if c is not None else "" for c in row)
            if row_text.strip(" |"):
                rows.append(row_text)
        if rows:
            sheets.append(f"[Sheet: {sheet_name}]\n" + "\n".join(rows))

    wb.close()
    return "\n\n".join(sheets), []  # no embedded images for XLSX
